"""PPTX parser — markitdown for text, python-pptx for bitmaps (v0.4.4).

Phase 0 benchmark (worklog/001) put markitdown at 23K text vs 170 for a
hand-rolled python-pptx parser. v0.3 ships text delegation; v0.4.4 closes
worklog/014 § 5-1 by adding direct python-pptx access for the bitmaps that
markitdown emits as ``![](filename)`` text placeholders (worklog/016).

Heavy optional dependency: ``markitdown[pptx]`` pulls in ``onnxruntime`` +
``numpy`` + ``magika`` (~80MB) and incidentally provides ``python-pptx``.
Therefore this parser is **opt-in only** — ``pip install korean-doc-parser[pptx]``.
When markitdown isn't installed, this module imports cleanly but doesn't
register itself with the registry, so a core-only install gets a friendly
:class:`UnsupportedFormatError` instead of an ``ImportError``.

The parser registers itself with the global registry on import **iff**
markitdown is importable.
"""

from __future__ import annotations

import hashlib
import tempfile
from collections.abc import Iterator
from pathlib import Path
from typing import TYPE_CHECKING

from korean_doc_parser.core import (
    BaseParser,
    ExtractedImage,
    ParseMetadata,
    ParseResult,
    register_parser,
)
from korean_doc_parser.exceptions import ParseError

if TYPE_CHECKING:
    from pptx.shapes.picture import Picture

__all__ = ["PptxParser"]


class PptxParser(BaseParser):
    """Parse ``.pptx`` files — text via markitdown, bitmaps via python-pptx.

    v0.4.4 contract:
    * ``markdown`` — markitdown's slide-by-slide output (``<!-- Slide number: N -->``
      delimiters preserved)
    * ``tables`` — empty. markitdown emits tables as markdown text inline with
      paragraphs; promoting them to ``ParsedTable`` requires re-parsing the
      markdown which loses fidelity. Deferred.
    * ``images`` — populated from python-pptx ``MSO_SHAPE_TYPE.PICTURE`` shapes
      (groups walked recursively). v0.3 returned ``[]`` here; v0.4.4 fills it.
      Notes:
      - ``page_no`` = 1-based slide index (notes / master / layout slides are
        skipped — only main slides)
      - ``bbox`` is in PPTX EMU units (1 inch = 914400 EMU); other parsers use
        pixels — downstream code should not assume a unified unit. Future
        normalisation is a v0.5+ Pipeline-side decision.
      - Non-standard image formats (``.wdp`` / HD Photo) keep ``width=height=0``
        but ``sha256`` + ``file_path`` are still valid for downstream Vision.
    * ``metadata.title`` — best-effort from python-pptx ``core_properties``.
    """

    @property
    def supported_extensions(self) -> tuple[str, ...]:
        return (".pptx",)

    def parse(self, path: Path) -> ParseResult:
        from markitdown import MarkItDown

        try:
            md = MarkItDown()
            result = md.convert(str(path))
            markdown = result.text_content or ""

            metadata = ParseMetadata(
                format="pptx",
                file_path=path,
                title=_extract_title(path),
            )

            return ParseResult(
                markdown=markdown,
                metadata=metadata,
                tables=[],
                images=_collect_pictures(path),
            )
        except ParseError:
            raise
        except Exception as exc:
            msg = f"Failed to parse PPTX {path}: {exc}"
            raise ParseError(msg) from exc


def _extract_title(path: Path) -> str | None:
    """Best-effort title from PPTX core properties.

    Returns ``None`` on any failure — title is informational, not contractual.
    python-pptx is pulled in by ``markitdown[pptx]`` extras so importing it is
    safe inside the gated branch.
    """
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        title = prs.core_properties.title
        return title or None
    except Exception:
        return None


# ─────────────────────────────────────────────────────────────────────────────
# Bitmap extraction (v0.4.4, worklog/016)
# ─────────────────────────────────────────────────────────────────────────────


def _collect_pictures(path: Path) -> list[ExtractedImage]:
    """Walk all slides and return one ``ExtractedImage`` per picture shape.

    A best-effort path: parser errors in python-pptx (corrupt PPTX, unknown
    shape type) fall through to ``[]`` rather than propagating, because the
    text-side ``markdown`` is the primary contract and we don't want bitmap
    extraction failure to fail the whole parse. The wrapping ``parse()`` still
    surfaces fatal PPTX errors from markitdown itself.
    """
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
    except Exception:
        return []
    out: list[ExtractedImage] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for order, shape in enumerate(_iter_picture_shapes(slide.shapes), start=1):
            extracted = _shape_to_image(shape, slide_idx, order)
            if extracted is not None:
                out.append(extracted)
    return out


def _iter_picture_shapes(shapes: object) -> Iterator[Picture]:
    """Yield ``PICTURE`` shapes from a shape collection, walking ``GROUP``s.

    The ``shapes`` argument is duck-typed — ``slide.shapes`` and a group
    shape's ``.shapes`` expose the same iteration protocol. python-pptx's
    enum import is deferred so the module loads in a core-only install.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:  # type: ignore[attr-defined]
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif st == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)


def _shape_to_image(shape: Picture, slide_no: int, order_in_slide: int) -> ExtractedImage | None:
    """Materialise one picture shape as an ``ExtractedImage`` + tempfile.

    Returns ``None`` only when ``shape.image`` itself raises (corrupt or
    placeholder picture without an embedded blob). All other failures (PIL
    can't decode ``.wdp`` etc.) keep the row with ``width=height=0`` —
    sha256 + file_path are still usable for downstream Vision.
    """
    try:
        image = shape.image
        blob: bytes = image.blob
        content_type: str = image.content_type or "application/octet-stream"
        ext: str = image.ext or "bin"
    except Exception:
        return None

    sha = hashlib.sha256(blob).hexdigest()
    width, height = _safe_image_size(blob)

    with tempfile.NamedTemporaryFile(
        prefix=f"kdp_pptx_s{slide_no}_{order_in_slide}_",
        suffix=f".{ext}",
        delete=False,
    ) as fh:
        fh.write(blob)
        file_path = fh.name

    return ExtractedImage(
        page_no=slide_no,
        section_no=None,
        bbox=_safe_bbox(shape),
        order_in_page=order_in_slide,
        text_before="",
        text_after="",
        section_title=None,
        file_path=file_path,
        sha256=sha,
        width=width,
        height=height,
        size_bytes=len(blob),
        mime_type=content_type,
        detected_caption=None,
        caption_method=None,
        caption_pattern_score=0.0,
    )


def _safe_image_size(data: bytes) -> tuple[int, int]:
    """PIL-based (width, height); returns (0, 0) for formats PIL can't open
    (e.g. ``.wdp`` / HD Photo, which python-pptx hands out as raw bytes)."""
    from io import BytesIO

    from PIL import Image

    try:
        with Image.open(BytesIO(data)) as img:
            return int(img.width), int(img.height)
    except Exception:
        return 0, 0


def _safe_bbox(shape: Picture) -> tuple[float, float, float, float] | None:
    """``(left, top, right, bottom)`` in EMU. ``None`` if any coord is missing.

    EMU (English Metric Units): 1 inch = 914400 EMU. PPTX shapes carry these
    natively; conversion to pixels needs a viewport DPI which the parser
    doesn't know. Downstream code that needs a unified unit should normalise.
    """
    try:
        left = float(shape.left) if shape.left is not None else None
        top = float(shape.top) if shape.top is not None else None
        width = float(shape.width) if shape.width is not None else None
        height = float(shape.height) if shape.height is not None else None
    except Exception:
        return None
    if None in (left, top, width, height):
        return None
    # mypy narrows after the None check — assert helps the strict checker.
    assert left is not None and top is not None and width is not None and height is not None
    return (left, top, left + width, top + height)


# ─────────────────────────────────────────────────────────────────────────────
# Conditional auto-register — only when markitdown is available
# ─────────────────────────────────────────────────────────────────────────────

try:
    import markitdown  # noqa: F401

    register_parser(PptxParser())
except ImportError:
    # Core-only install — .pptx will raise UnsupportedFormatError, which is
    # the documented behaviour. Hint the user in the message at the caller.
    pass

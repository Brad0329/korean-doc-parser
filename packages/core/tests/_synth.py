"""Synthetic fixture generators — code-built tiny documents for tests.

Each generator returns a ``Path`` to a freshly-built file inside the caller-
provided directory. Generators stay format-isolated so a missing optional dep
(e.g. ``python-docx``) only breaks the formats that need it.

Used by :mod:`conftest` session-scoped pytest fixtures so the cost is paid
once per test run.

Ground-truth JSON is written next to the fixture so per-parser tests can call
:func:`korean_doc_parser.tests._gt.load_ground_truth` uniformly.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

__all__ = [
    "build_docx_simple",
    "build_docx_with_image",
    "build_docx_with_table",
    "build_hwpx_simple",
    "build_hwpx_with_image",
    "build_hwpx_with_image_corrupt",
    "build_hwpx_with_image_resources",
    "build_hwpx_with_table",
    "build_pdf_multipage",
    "build_pdf_simple",
    "build_pdf_with_image",
    "build_pdf_with_image_cmyk",
    "build_pdf_with_image_jpeg",
    "build_pdf_with_multi_images",
    "build_pdf_with_table",
    "build_pptx_multislide",
    "build_pptx_simple",
    "build_pptx_with_image",
    "build_pptx_with_table",
]


def _write_gt(fixture_path: Path, gt: dict[str, Any]) -> None:
    """Persist ground truth JSON alongside ``fixture_path``."""
    gt_path = fixture_path.with_name(fixture_path.name + ".gt.json")
    gt_path.write_text(
        json.dumps(gt, ensure_ascii=False, indent=2, sort_keys=True),
        encoding="utf-8",
    )


# ─────────────────────────────────────────────────────────────────────────────
# DOCX
# ─────────────────────────────────────────────────────────────────────────────


def build_docx_simple(dest_dir: Path) -> Path:
    """Build a 1-paragraph DOCX with a heading. Returns the fixture path."""
    from docx import Document

    path = dest_dir / "docx_simple.docx"
    doc = Document()
    doc.add_heading("테스트 문서 제목", level=1)
    doc.add_paragraph(
        "이것은 합성 테스트용 한국어 DOCX 문서입니다. "
        "본문 길이를 검증하기 위한 충분한 길이의 텍스트를 포함합니다."
    )
    doc.save(str(path))

    _write_gt(
        path,
        {
            "expected_format": "docx",
            "expected_table_count": 0,
            "expected_image_count": 0,
            "expected_sections": ["테스트 문서 제목"],
            "expected_keywords": ["합성 테스트용", "한국어 DOCX"],
            "expected_text_length_range": [40, 500],
        },
    )
    return path


def build_docx_with_table(dest_dir: Path) -> Path:
    """Build a DOCX containing a single 3-by-3 table. Returns the fixture path."""
    from docx import Document

    path = dest_dir / "docx_with_table.docx"
    doc = Document()
    doc.add_heading("표 포함 문서", level=1)
    doc.add_paragraph("아래 표는 예시 데이터입니다.")

    table = doc.add_table(rows=3, cols=3)
    headers = ["항목", "값", "비고"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    table.rows[1].cells[0].text = "매출"
    table.rows[1].cells[1].text = "1,000"
    table.rows[1].cells[2].text = "단위: 백만 원"
    table.rows[2].cells[0].text = "비용"
    table.rows[2].cells[1].text = "600"
    table.rows[2].cells[2].text = "고정비 포함"

    doc.save(str(path))

    _write_gt(
        path,
        {
            "expected_format": "docx",
            "expected_table_count": 1,
            "expected_image_count": 0,
            "expected_sections": ["표 포함 문서"],
            # Note: table cells live in result.tables (not in markdown),
            # so the keyword check targets paragraph body only.
            "expected_keywords": ["예시 데이터"],
            "expected_text_length_range": [20, 500],
        },
    )
    return path


def build_docx_with_image(dest_dir: Path) -> Path:
    """Synthetic DOCX containing one inline image (32-by-32 PNG)."""
    from io import BytesIO

    from docx import Document
    from PIL import Image

    path = dest_dir / "docx_with_image.docx"
    doc = Document()
    doc.add_heading("이미지 포함 문서", level=1)
    doc.add_paragraph("아래 그림은 임베디드 이미지입니다.")

    png_buf = BytesIO()
    Image.new("RGB", (32, 32), (0, 128, 0)).save(png_buf, format="PNG")
    png_buf.seek(0)
    doc.add_picture(png_buf, width=None)
    doc.save(str(path))

    _write_gt(
        path,
        {
            "expected_format": "docx",
            "expected_table_count": 0,
            "expected_image_count": 1,
            "expected_sections": ["이미지 포함 문서"],
            "expected_keywords": ["임베디드 이미지"],
            "expected_text_length_range": [10, 500],
        },
    )
    return path


# ─────────────────────────────────────────────────────────────────────────────
# PDF — reportlab. English text only (default fonts lack Korean glyphs);
# Korean PDF coverage comes from real `samples/*.pdf` fixtures, not synth.
# ─────────────────────────────────────────────────────────────────────────────


def build_pdf_simple(dest_dir: Path) -> Path:
    """Synthetic single-page PDF with English text."""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    path = dest_dir / "pdf_simple.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setTitle("Synthetic PDF Title")
    c.setAuthor("doc_parser tests")
    c.drawString(72, 800, "Synthetic PDF Fixture")
    c.drawString(72, 770, "This is a minimal one-page PDF for parser tests.")
    c.showPage()
    c.save()

    _write_gt(
        path,
        {
            "expected_format": "pdf",
            "expected_page_count": 1,
            "expected_table_count": 0,
            "expected_image_count": 0,
            "expected_sections": ["Synthetic PDF Fixture"],
            "expected_keywords": ["minimal one-page PDF"],
            "expected_text_length_range": [30, 500],
            "expected_title": "Synthetic PDF Title",
        },
    )
    return path


def build_pdf_multipage(dest_dir: Path) -> Path:
    """Synthetic 3-page PDF — exercises page_count + multi-page text join."""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    path = dest_dir / "pdf_multipage.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    for i in range(3):
        c.drawString(72, 800, f"Page {i + 1} of 3")
        c.drawString(72, 770, f"Content on page {i + 1}.")
        c.showPage()
    c.save()

    _write_gt(
        path,
        {
            "expected_format": "pdf",
            "expected_page_count": 3,
            "expected_table_count": 0,
            "expected_image_count": 0,
            "expected_sections": ["Page 1 of 3", "Page 2 of 3", "Page 3 of 3"],
            "expected_text_length_range": [40, 500],
        },
    )
    return path


def build_pdf_with_table(dest_dir: Path) -> Path:
    """Synthetic single-page PDF with one 3-by-3 table."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table

    path = dest_dir / "pdf_with_table.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=A4)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("PDF with Table Fixture", styles["Heading1"]),
        Spacer(1, 12),
        Table(
            [
                ["Item", "Value", "Note"],
                ["Revenue", "1000", "Unit: million KRW"],
                ["Cost", "600", "Includes fixed cost"],
            ],
            style=[
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ],
        ),
    ]
    doc.build(story)

    _write_gt(
        path,
        {
            "expected_format": "pdf",
            "expected_page_count": 1,
            "expected_table_count": 1,
            "expected_image_count": 0,
            "expected_sections": ["PDF with Table Fixture"],
            "expected_keywords": ["Revenue", "Cost"],
            "expected_text_length_range": [50, 1500],
        },
    )
    return path


def build_pdf_with_image(dest_dir: Path) -> Path:
    """Synthetic single-page PDF embedding a 32-by-32 red square PNG."""
    from io import BytesIO

    from PIL import Image
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    path = dest_dir / "pdf_with_image.pdf"

    png_buf = BytesIO()
    Image.new("RGB", (32, 32), (255, 0, 0)).save(png_buf, format="PNG")
    png_buf.seek(0)

    c = canvas.Canvas(str(path), pagesize=A4)
    c.drawString(72, 800, "PDF with Image Fixture")
    c.drawImage(ImageReader(png_buf), 100, 700, width=32, height=32)
    c.showPage()
    c.save()

    _write_gt(
        path,
        {
            "expected_format": "pdf",
            "expected_page_count": 1,
            "expected_table_count": 0,
            "expected_image_count": 1,
            "expected_sections": ["PDF with Image Fixture"],
            "expected_text_length_range": [10, 500],
        },
    )
    return path


def build_pdf_with_image_jpeg(dest_dir: Path) -> Path:
    """Single-page PDF embedding a 48x32 JPEG (RGB).

    Exercises pypdf's ``/DCTDecode`` filter path. Different dimensions from the
    PNG fixture so size assertions can distinguish them in test failures.
    """
    from io import BytesIO

    from PIL import Image
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    path = dest_dir / "pdf_with_image_jpeg.pdf"

    jpg_buf = BytesIO()
    Image.new("RGB", (48, 32), (0, 200, 100)).save(jpg_buf, format="JPEG", quality=80)
    jpg_buf.seek(0)

    c = canvas.Canvas(str(path), pagesize=A4)
    c.drawString(72, 800, "PDF with JPEG Image")
    c.drawImage(ImageReader(jpg_buf), 100, 700, width=48, height=32)
    c.showPage()
    c.save()

    _write_gt(
        path,
        {
            "expected_format": "pdf",
            "expected_page_count": 1,
            "expected_table_count": 0,
            "expected_image_count": 1,
            "expected_sections": ["PDF with JPEG Image"],
            "expected_text_length_range": [10, 500],
        },
    )
    return path


def build_pdf_with_image_cmyk(dest_dir: Path) -> Path:
    """Single-page PDF embedding a 40x40 CMYK JPEG.

    Exercises the CMYK colorspace branch of pypdf's decoder. Scanned/print-ready
    PDFs (common in 입찰공고) frequently embed CMYK JPEGs.
    """
    from io import BytesIO

    from PIL import Image
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    path = dest_dir / "pdf_with_image_cmyk.pdf"

    cmyk_buf = BytesIO()
    # CMYK with C=50% M=20% Y=70% K=10% — non-trivial colorspace
    Image.new("CMYK", (40, 40), (128, 51, 178, 26)).save(cmyk_buf, format="JPEG", quality=80)
    cmyk_buf.seek(0)

    c = canvas.Canvas(str(path), pagesize=A4)
    c.drawString(72, 800, "PDF with CMYK JPEG")
    c.drawImage(ImageReader(cmyk_buf), 100, 700, width=40, height=40)
    c.showPage()
    c.save()

    _write_gt(
        path,
        {
            "expected_format": "pdf",
            "expected_page_count": 1,
            "expected_table_count": 0,
            "expected_image_count": 1,
            "expected_sections": ["PDF with CMYK JPEG"],
            "expected_text_length_range": [10, 500],
        },
    )
    return path


def build_pdf_with_multi_images(dest_dir: Path) -> Path:
    """Two-page PDF with one image per page — exercises per-page index matching.

    Page 1: 24x24 blue PNG, Page 2: 24x24 yellow PNG. Verifies that pdfplumber
    bbox / pypdf bitmap pairing stays correct across page boundaries.
    """
    from io import BytesIO

    from PIL import Image
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    path = dest_dir / "pdf_with_multi_images.pdf"

    blue_buf = BytesIO()
    Image.new("RGB", (24, 24), (0, 0, 255)).save(blue_buf, format="PNG")
    blue_buf.seek(0)

    yellow_buf = BytesIO()
    Image.new("RGB", (24, 24), (255, 255, 0)).save(yellow_buf, format="PNG")
    yellow_buf.seek(0)

    c = canvas.Canvas(str(path), pagesize=A4)
    c.drawString(72, 800, "Page 1 with Blue PNG")
    c.drawImage(ImageReader(blue_buf), 100, 700, width=24, height=24)
    c.showPage()
    c.drawString(72, 800, "Page 2 with Yellow PNG")
    c.drawImage(ImageReader(yellow_buf), 100, 700, width=24, height=24)
    c.showPage()
    c.save()

    _write_gt(
        path,
        {
            "expected_format": "pdf",
            "expected_page_count": 2,
            "expected_table_count": 0,
            "expected_image_count": 2,
            "expected_sections": ["Page 1 with Blue PNG", "Page 2 with Yellow PNG"],
            "expected_text_length_range": [20, 500],
        },
    )
    return path


# ─────────────────────────────────────────────────────────────────────────────
# HWPX — direct OOXML-style XML (no external dependency)
# ─────────────────────────────────────────────────────────────────────────────

_HWPX_MIMETYPE = "application/hwp+zip"

_HWPX_CONTENT_HPF = """<?xml version="1.0" encoding="UTF-8"?>
<opf:package xmlns:opf="http://www.idpf.org/2007/opf">
  <opf:metadata>
    <opf:title>HWPX 합성 픽스처</opf:title>
  </opf:metadata>
</opf:package>
"""

_HWPX_SECTION0_XML = """<?xml version="1.0" encoding="UTF-8"?>
<hp:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p>
    <hp:run><hp:t>HWPX 합성 픽스처</hp:t></hp:run>
  </hp:p>
  <hp:p>
    <hp:run><hp:t>이것은 코드로 생성한 최소 HWPX 본문입니다.</hp:t></hp:run>
  </hp:p>
</hp:sec>
"""


_HWPX_SECTION0_TABLE_XML = """<?xml version="1.0" encoding="UTF-8"?>
<hp:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p>
    <hp:run><hp:t>표 포함 HWPX</hp:t></hp:run>
  </hp:p>
  <hp:p>
    <hp:tbl>
      <hp:tr>
        <hp:tc><hp:run><hp:t>A</hp:t></hp:run></hp:tc>
        <hp:tc><hp:run><hp:t>B</hp:t></hp:run></hp:tc>
      </hp:tr>
      <hp:tr>
        <hp:tc><hp:run><hp:t>1</hp:t></hp:run></hp:tc>
        <hp:tc><hp:run><hp:t>2</hp:t></hp:run></hp:tc>
      </hp:tr>
    </hp:tbl>
  </hp:p>
</hp:sec>
"""


def build_hwpx_simple(dest_dir: Path) -> Path:
    """Build a minimal valid HWPX (ZIP with mimetype + content + section).

    The XML is intentionally simple so the parser test focuses on structure
    handling, not Hancom's full schema. Real HWPX from 한/글 has much richer
    namespacing; the parser must tolerate this minimal shape too.
    """
    import zipfile

    path = dest_dir / "hwpx_simple.hwpx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        # mimetype must be stored (not deflated) and first per the OPF spec;
        # for synthetic-test purposes we don't enforce STORED here.
        zf.writestr("mimetype", _HWPX_MIMETYPE)
        zf.writestr("Contents/content.hpf", _HWPX_CONTENT_HPF)
        zf.writestr("Contents/section0.xml", _HWPX_SECTION0_XML)

    _write_gt(
        path,
        {
            "expected_format": "hwpx",
            "expected_table_count": 0,
            "expected_image_count": 0,
            "expected_sections": ["HWPX 합성 픽스처"],
            "expected_keywords": ["코드로 생성한 최소 HWPX"],
            "expected_text_length_range": [20, 500],
            "expected_title": "HWPX 합성 픽스처",
        },
    )
    return path


def build_hwpx_with_table(dest_dir: Path) -> Path:
    """Build a minimal HWPX containing one 2-by-2 table inside section0."""
    import zipfile

    path = dest_dir / "hwpx_with_table.hwpx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", _HWPX_MIMETYPE)
        zf.writestr("Contents/content.hpf", _HWPX_CONTENT_HPF)
        zf.writestr("Contents/section0.xml", _HWPX_SECTION0_TABLE_XML)

    _write_gt(
        path,
        {
            "expected_format": "hwpx",
            "expected_table_count": 1,
            "expected_image_count": 0,
            "expected_sections": ["표 포함 HWPX"],
            "expected_text_length_range": [5, 200],
        },
    )
    return path


# ── HWPX image fixtures (exercise _extract_images + _inspect_image branches) ──

_HWPX_SECTION0_IMAGE_XML = """<?xml version="1.0" encoding="UTF-8"?>
<hp:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p>
    <hp:run><hp:t>이미지 포함 HWPX</hp:t></hp:run>
  </hp:p>
  <hp:p>
    <hp:run><hp:t>아래는 BinData 경로의 임베디드 이미지입니다.</hp:t></hp:run>
  </hp:p>
</hp:sec>
"""

_HWPX_SECTION0_IMAGE_RESOURCES_XML = """<?xml version="1.0" encoding="UTF-8"?>
<hp:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p>
    <hp:run><hp:t>Resources 경로 이미지 HWPX</hp:t></hp:run>
  </hp:p>
</hp:sec>
"""

_HWPX_SECTION0_IMAGE_CORRUPT_XML = """<?xml version="1.0" encoding="UTF-8"?>
<hp:sec xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p>
    <hp:run><hp:t>손상 이미지 HWPX</hp:t></hp:run>
  </hp:p>
</hp:sec>
"""


def build_hwpx_with_image(dest_dir: Path) -> Path:
    """HWPX with a 32x32 PNG under ``BinData/``.

    Exercises: BinData/ path branch, PIL-success branch of ``_inspect_image``,
    PNG mime resolution, ``.png`` suffix path of the tempfile writer.
    """
    import zipfile
    from io import BytesIO

    from PIL import Image

    path = dest_dir / "hwpx_with_image.hwpx"
    png_buf = BytesIO()
    Image.new("RGB", (32, 32), (0, 0, 255)).save(png_buf, format="PNG")

    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", _HWPX_MIMETYPE)
        zf.writestr("Contents/content.hpf", _HWPX_CONTENT_HPF)
        zf.writestr("Contents/section0.xml", _HWPX_SECTION0_IMAGE_XML)
        zf.writestr("BinData/image1.png", png_buf.getvalue())

    _write_gt(
        path,
        {
            "expected_format": "hwpx",
            "expected_table_count": 0,
            "expected_image_count": 1,
            "expected_sections": ["이미지 포함 HWPX"],
            "expected_keywords": ["BinData 경로"],
            "expected_text_length_range": [10, 500],
        },
    )
    return path


def build_hwpx_with_image_resources(dest_dir: Path) -> Path:
    """HWPX with a 48x24 JPEG under ``Contents/Resources/``.

    Exercises: Contents/Resources/ path branch, JPEG format detection,
    ``.jpg`` suffix path. Pairs with ``build_hwpx_with_image`` to cover both
    media roots that ``_extract_images`` recognizes.
    """
    import zipfile
    from io import BytesIO

    from PIL import Image

    path = dest_dir / "hwpx_with_image_resources.hwpx"
    jpg_buf = BytesIO()
    Image.new("RGB", (48, 24), (200, 100, 0)).save(jpg_buf, format="JPEG")

    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", _HWPX_MIMETYPE)
        zf.writestr("Contents/content.hpf", _HWPX_CONTENT_HPF)
        zf.writestr("Contents/section0.xml", _HWPX_SECTION0_IMAGE_RESOURCES_XML)
        zf.writestr("Contents/Resources/image1.jpg", jpg_buf.getvalue())

    _write_gt(
        path,
        {
            "expected_format": "hwpx",
            "expected_table_count": 0,
            "expected_image_count": 1,
            "expected_sections": ["Resources 경로 이미지 HWPX"],
            "expected_text_length_range": [5, 200],
        },
    )
    return path


def build_hwpx_with_image_corrupt(dest_dir: Path) -> Path:
    """HWPX with non-image bytes under ``BinData/broken.bin``.

    Exercises ``_inspect_image``'s exception path: PIL fails to open the bytes,
    parser must fall back to ``(0, 0, application/octet-stream)`` rather than
    raise. Real-world payoff: corrupted media inside otherwise-valid HWPX
    should not crash the whole document parse.
    """
    import zipfile

    path = dest_dir / "hwpx_with_image_corrupt.hwpx"

    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", _HWPX_MIMETYPE)
        zf.writestr("Contents/content.hpf", _HWPX_CONTENT_HPF)
        zf.writestr("Contents/section0.xml", _HWPX_SECTION0_IMAGE_CORRUPT_XML)
        zf.writestr("BinData/broken.bin", b"not a real image, this is junk bytes")

    _write_gt(
        path,
        {
            "expected_format": "hwpx",
            "expected_table_count": 0,
            "expected_image_count": 1,
            "expected_sections": ["손상 이미지 HWPX"],
            "expected_text_length_range": [5, 200],
        },
    )
    return path


# ─────────────────────────────────────────────────────────────────────────────
# PPTX — python-pptx synthesis. Korean fixtures verify markitdown's slide-by-
# slide unicode handling; an image fixture proves the placeholder pattern.
# Real-world Korean PPTX coverage comes from `samples/*.pptx` (skip-guarded).
# ─────────────────────────────────────────────────────────────────────────────


def build_pptx_simple(dest_dir: Path) -> Path:
    """1-slide PPTX with a Korean title + body. Quickest regression target."""
    from pptx import Presentation
    from pptx.util import Inches

    path = dest_dir / "pptx_simple.pptx"
    prs = Presentation()
    prs.core_properties.title = "합성 PPTX 단일 슬라이드"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "합성 PPTX 단일 슬라이드"
    body = slide.placeholders[1]
    body.text = "이것은 합성 테스트용 한국어 PPTX 본문입니다."
    body.text_frame.add_paragraph().text = "두 번째 단락."
    _ = Inches  # silence unused-import; keeps the helper import stable
    prs.save(str(path))

    _write_gt(
        path,
        {
            "expected_format": "pptx",
            "expected_table_count": 0,
            "expected_image_count": 0,
            "expected_keywords": ["합성", "한국어 PPTX"],
            "expected_text_length_range": [40, 1500],
            "expected_title": "합성 PPTX 단일 슬라이드",
        },
    )
    return path


def build_pptx_multislide(dest_dir: Path) -> Path:
    """3-slide PPTX — exercises markitdown's per-slide section markers."""
    from pptx import Presentation

    path = dest_dir / "pptx_multislide.pptx"
    prs = Presentation()
    for i in range(1, 4):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"슬라이드 {i}"
        slide.placeholders[1].text = f"이 슬라이드의 본문은 {i}번입니다."
    prs.save(str(path))

    _write_gt(
        path,
        {
            "expected_format": "pptx",
            "expected_table_count": 0,
            "expected_image_count": 0,
            "expected_keywords": ["슬라이드 1", "슬라이드 2", "슬라이드 3"],
            "expected_text_length_range": [60, 2000],
        },
    )
    return path


def build_pptx_with_table(dest_dir: Path) -> Path:
    """PPTX with a single 3-row x 3-column table on slide 1."""
    from pptx import Presentation
    from pptx.util import Inches

    path = dest_dir / "pptx_with_table.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title + blank
    slide.shapes.title.text = "표 포함 PPTX"

    shape = slide.shapes.add_table(
        rows=3,
        cols=3,
        left=Inches(1),
        top=Inches(2),
        width=Inches(6),
        height=Inches(2),
    )
    tbl = shape.table
    headers = ["항목", "값", "비고"]
    for i, h in enumerate(headers):
        tbl.cell(0, i).text = h
    tbl.cell(1, 0).text = "매출"
    tbl.cell(1, 1).text = "1,000"
    tbl.cell(1, 2).text = "단위: 백만 원"
    tbl.cell(2, 0).text = "비용"
    tbl.cell(2, 1).text = "600"
    tbl.cell(2, 2).text = "고정비 포함"

    prs.save(str(path))

    _write_gt(
        path,
        {
            "expected_format": "pptx",
            "expected_table_count": 0,  # v0.3 contract: tables stay inline in markdown
            "expected_image_count": 0,
            "expected_keywords": ["표 포함 PPTX", "매출", "비용"],
            "expected_text_length_range": [30, 2000],
        },
    )
    return path


def build_pptx_with_image(dest_dir: Path) -> Path:
    """PPTX embedding a 64x64 red square PNG on slide 1."""
    from io import BytesIO

    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    path = dest_dir / "pptx_with_image.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "이미지 포함 PPTX"

    png_buf = BytesIO()
    Image.new("RGB", (64, 64), (255, 0, 0)).save(png_buf, format="PNG")
    png_buf.seek(0)
    slide.shapes.add_picture(png_buf, Inches(1), Inches(2), Inches(1), Inches(1))

    prs.save(str(path))

    _write_gt(
        path,
        {
            "expected_format": "pptx",
            "expected_table_count": 0,
            "expected_image_count": 0,  # v0.3: bitmap promotion deferred
            "expected_keywords": ["이미지 포함 PPTX"],
            "expected_text_length_range": [10, 2000],
        },
    )
    return path
